from flask import Flask, request
from groq import Groq
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from pptx.dml.color import RGBColor

app = Flask(__name__)

@app.route('/api/generate', methods=['POST'])
def generate():
    print("came")
    if request.method == 'POST':
        print("came")
        data = request.json 
        query = data.get('query')
        client = Groq(api_key="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")
        completion = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "generate / rectify grammar - give a one line response - neither too long nor too short (maximum 5 words) - just the title"
                },
                {
                    "role": "user",
                    "content": f"query[{query}]"
                }
            ],
            temperature=1,
            max_tokens=8192,
            top_p=1,
            stream=True,
            stop=None,
        )
        query_result = ""
        for chunk in completion:
            query_result += chunk.choices[0].delta.content or ""
        print(query_result)

        slides = 5

        completion = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": f"generate {slides} ppt slide titles in order (keep them sweet and short) based on the query[] - split the titles with delimiter '%%%' and return all titles within a codeblock ``` - do not mention slide count - do not include asterisks * - dont generate question tags - max words - (6 to 8) - dont include double quotes ,just plain text"
                },
                {
                    "role": "user",
                    "content": f"query[{query_result}]"
                }
            ],
            temperature=1,
            max_tokens=8192,
            top_p=1,
            stream=True,
            stop=None,
        )
        slide_points_full = ""
        for chunk in completion:
            slide_points_full += chunk.choices[0].delta.content or ""

        slide_points = slide_points_full[slide_points_full.find('```')+4:slide_points_full.rfind('```')-1].replace('\n',"").split("%%%")

        slide_content = []
        for point in slide_points:
            if point.strip() != "":
                completion = client.chat.completions.create(
                    model="llama3-8b-8192",
                    messages=[
                        {
                            "role": "system",
                            "content": f"generate a paragraph (50-60 words STRICTLY) for slide content ONLY (neither too big nor too small) for title[] - return the body only without any unwanted text **do not use linebreak** - do not use asterisks* - use justify text align - do not mention anything else, give only the slide content paragraph as a response - return only information related to the point[] not statements like 'here is..','some of ...' etc - do not use linebreak generate entire response in one line - dont use asterisks * "
                        },
                        {
                            "role": "user",
                            "content": f"[{point}]"
                        }
                    ],
                    temperature=1,
                    max_tokens=8192,
                    top_p=1,
                    stream=True,
                    stop=None,
                )
                temp_slide_content = ""
                for chunk in completion:
                    temp_slide_content += chunk.choices[0].delta.content or ""
                slide_content.append(temp_slide_content)

        slide_width = Inches(10)  
        slide_height = Inches(5.625) 

        prs = Presentation()
        prs.slide_width = slide_width
        prs.slide_height = slide_height

        slide = prs.slides.add_slide(prs.slide_layouts[5])  

        img_path = '/home/sid/Desktop/hub/PPT-AI/server/styles/template1/1.png'
        img = Image.open(img_path)
        img_width, img_height = img.size

        max_width = slide_width  
        max_height = slide_height 
        scale_width = max_width / img_width
        scale_height = max_height / img_height
        scale = min(scale_width, scale_height)

        left = (slide_width - img_width * scale) / 2
        top = (slide_height - img_height * scale) / 2

        pic = slide.shapes.add_picture(img_path, left, top, width=img_width * scale, height=img_height * scale)
        title_width = Inches(8) 
        title_left = (slide_width - title_width) / 2
        title = slide.shapes.add_textbox(title_left, Inches(2), title_width, Inches(1))
        title.text = query_result
        title.text_frame.word_wrap = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
        title.text_frame.paragraphs[0].font.bold = True  
        title.text_frame.paragraphs[0].font.size = Inches(0.4)  
        title.text_frame.paragraphs[0].font.name = 'Courier New' 
        title.text_frame.paragraphs[0].alignment = 2

        img_path2 = '/home/sid/Desktop/hub/PPT-AI/server/styles/template1/2.png'
        for count in range(len(slide_points)):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  

            left = (slide_width - img_width * scale) / 2
            top = (slide_height - img_height * scale) / 2

            pic = slide.shapes.add_picture(img_path2, left, top, width=img_width * scale, height=img_height * scale)

            title = slide.shapes.add_textbox(left, Inches(0.2), img_width * scale, Inches(1))
            title.text = slide_points[count]
            title.text_frame.word_wrap = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.size = Inches(0.3) 
            title.text_frame.paragraphs[0].font.name = 'Courier New' 
            title.text_frame.paragraphs[0].alignment = 2
            title = slide.shapes.add_textbox(left+Inches(0.25), Inches(1.5), Inches(4), Inches(1))
            title.text = slide_content[count]
            title.text_frame.word_wrap = True  
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
            title.text_frame.paragraphs[0].font.bold = False
            title.text_frame.paragraphs[0].font.size = Inches(0.15) 
            title.text_frame.paragraphs[0].font.name = 'Courier New' 
            title.text_frame.paragraphs[0].alignment = 2

        prs.save('slide_genie.pptx')

        return 'Presentation generated successfully!'
    else:
        return 'Method Not Allowed', 405

if __name__ == '__main__':
    app.run(debug=True)
