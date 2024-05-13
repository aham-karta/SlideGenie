from flask import Flask, request
from flask_cors import CORS
from groq import Groq
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from pptx.dml.color import RGBColor

app = Flask(__name__)
CORS(app)

@app.route('/api/generate', methods=['POST'])
def generate():
    if request.method == 'POST':
        print("came")
        data = request.json 
        query = data.get('query')
        slides=int(data.get('slides'))
        if(slides<=0 or slides>30):
            return "number of slides <=0 or > 30"
        client = Groq(api_key="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")
        completion = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "generate / rectify grammar (if needed) - give a one line ppt title response - neither too long nor too short (maximum 5 words) - just the title - do not use asterisks* and double quotes"
                },
                {
                    "role": "user",
                    "content": f"query[{query}]"
                }
            ],
            temperature=1,
            max_tokens=8192,
            top_p=1,
            stream=True,
            stop=None,
        )
        query_result = ""
        for chunk in completion:
            query_result += chunk.choices[0].delta.content or ""
        print("made ppt title")

        completion = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": f"""generate {slides} ppt slide titles in order (keep them sweet and short) based on the query[] - split the titles with delimiter '%%' and return all titles within a codeblock ``` - do not mention slide count - do not include asterisks * - dont generate question tags - max words - (6 to 8) - dont include double quotes ,just plain text - sample response [title %% title %% title %% ...]"""
                },
                {
                    "role": "user",
                    "content": f"query[{query_result}], sample response [title1 %% title2 %% title3 %% ... titlen] , return the response only without any extra unecesary statements , dont respond like a conversation just respond like a code generator"
                }
            ],
            temperature=1,
            max_tokens=8192,
            top_p=1,
            stream=True,
            stop=None,
        )
        slide_points_full = ""
        for chunk in completion:
            slide_points_full += chunk.choices[0].delta.content or ""

        slide_points = slide_points_full[slide_points_full.find('```')+1:].replace('\n',"").split("%%")
        print("made slide titles")
        slide_content = []
        for point in slide_points:
            if point.strip() != "":
                completion = client.chat.completions.create(
                    model="llama3-8b-8192",
                    messages=[
                        {
                            "role": "system",
                            "content": f"generate a paragraph (50-60 words STRICTLY) for slide content ONLY (neither too big nor too small) for title[] - return the body only without any unwanted text **do not use linebreak** - do not use asterisks* - use justify text align - do not mention anything else, give only the slide content paragraph as a response - return only information related to the point[] not statements like 'here is..','some of ...' etc - do not use linebreak generate entire response in one line - dont use asterisks * "
                        },
                        {
                            "role": "user",
                            "content": f"[{point}]"
                        }
                    ],
                    temperature=1,
                    max_tokens=8192,
                    top_p=1,
                    stream=True,
                    stop=None,
                )
                temp_slide_content = ""
                for chunk in completion:
                    temp_slide_content += chunk.choices[0].delta.content or ""
                slide_content.append(temp_slide_content)
        print("generated slide content")
        slides=len(slide_content)
        print(len(slide_content))
        print(len(slide_points))
        slide_width = Inches(10)  
        slide_height = Inches(5.625) 

        prs = Presentation()
        prs.slide_width = slide_width
        prs.slide_height = slide_height
        slide_layout = prs.slide_layouts[0]
        for i in range(slides):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_points[i] if i < len(slide_points) else ""
            img_path = '/home/sid/Desktop/hub/PPT-AI/server/styles/template1/1.png'
            img = Image.open(img_path)
            img_width, img_height = img.size

            max_width = slide_width  
            max_height = slide_height 
            scale_width = max_width / img_width
            scale_height = max_height / img_height
            scale = min(scale_width, scale_height)

            left = (slide_width - img_width * scale) / 2
            top = (slide_height - img_height * scale) / 2

            pic = slide.shapes.add_picture(img_path, left, top, width=img_width * scale, height=img_height * scale)
            title_width = Inches(8) 
            title_left = (slide_width - title_width) / 2
            title = slide.shapes.add_textbox(title_left, Inches(2), title_width, Inches(1))
            title.text = query_result
            title.text_frame.word_wrap = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
            title.text_frame.paragraphs[0].font.bold = True  
            title.text_frame.paragraphs[0].font.size = Inches(0.4)  
            title.text_frame.paragraphs[0].font.name = 'Courier New' 
            title.text_frame.paragraphs[0].alignment = 2

            img_path2 = '/home/sid/Desktop/hub/PPT-AI/server/styles/template1/2.png'

            left = (slide_width - img_width * scale) / 2
            top = (slide_height - img_height * scale) / 2

            pic = slide.shapes.add_picture(img_path2, left, top, width=img_width * scale, height=img_height * scale)

            title = slide.shapes.add_textbox(left, Inches(0.2), img_width * scale, Inches(1))
            print(i)
            title.text = slide_points[i]
            title.text_frame.word_wrap = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.size = Inches(0.3) 
            title.text_frame.paragraphs[0].font.name = 'Courier New' 
            title.text_frame.paragraphs[0].alignment = 2
            title = slide.shapes.add_textbox(left+Inches(0.25), Inches(1.5), Inches(4), Inches(1))
            title.text = slide_content[i]
            title.text_frame.word_wrap = True  
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
            title.text_frame.paragraphs[0].font.bold = False
            title.text_frame.paragraphs[0].font.size = Inches(0.15) 
            title.text_frame.paragraphs[0].font.name = 'Courier New' 
            title.text_frame.paragraphs[0].alignment = 2
        prs.save('slide_genie.pptx')
        
        return 'Presentation generated successfully!'
    
    else:
        return 'Method Not Allowed', 405

if __name__ == '__main__':
    app.run(debug=True)